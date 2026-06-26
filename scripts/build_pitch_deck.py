#!/usr/bin/env python3
"""
Build Aniskwela-Grant-Pitch.pptx with earthen branding and embedded speaker notes.

Usage (from repo root):
    python scripts/build_pitch_deck.py
"""

from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUTPUT = REPO_ROOT / "Aniskwela-Grant-Pitch.pptx"

# DSD tokens
BG = RGBColor(0xFB, 0xF8, 0xF2)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE7, 0xE0, 0xD3)
PRIMARY = RGBColor(0x4B, 0x6B, 0xF1)
GROWTH = RGBColor(0x3F, 0x8E, 0x5B)
SOIL = RGBColor(0x6B, 0x4F, 0x3A)
TEXT = RGBColor(0x24, 0x1F, 0x1A)
MUTED = RGBColor(0x6E, 0x66, 0x5B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARNING = RGBColor(0xC8, 0x85, 0x2A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FOOTER = "ANISKWELA  ·  Axon Enjin"
DECK_LABEL = "Grant partner deck"
TOTAL_SLIDES = 11


def rgb_hex(c: RGBColor) -> str:
    return f"#{c}"


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_woven_accent(slide, left=0, top=0, width=Inches(13.333), height=Inches(0.08)) -> None:
    """Subtle top band suggesting banig/weave structure."""
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GROWTH
    band.line.fill.background()
    band.rotation = 0


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
    font_name: str = "Segoe UI",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color


def add_footer(slide, index: int, total: int) -> None:
    add_textbox(
        slide,
        Inches(0.6),
        Inches(7.05),
        Inches(5),
        Inches(0.35),
        FOOTER,
        size=10,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(10.5),
        Inches(7.05),
        Inches(2.3),
        Inches(0.35),
        f"{DECK_LABEL}  ·  {index}/{total}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_kicker(slide, text: str, top=Inches(0.55)) -> None:
    add_textbox(slide, Inches(0.75), top, Inches(4), Inches(0.35), text, size=11, bold=True, color=GROWTH)


def add_headline(slide, text: str, top=Inches(0.95)) -> None:
    add_textbox(slide, Inches(0.75), top, Inches(11.5), Inches(0.9), text, size=32, bold=True, color=SOIL)


def add_subhead(slide, text: str, top=Inches(1.75)) -> None:
    add_textbox(slide, Inches(0.75), top, Inches(11.5), Inches(0.7), text, size=17, color=MUTED)


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body: str,
    *,
    tag: str | None = None,
    tag_color: RGBColor = GROWTH,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = SURFACE
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)

    if tag:
        p0 = tf.paragraphs[0]
        p0.text = tag
        p0.space_after = Pt(4)
        r0 = p0.runs[0]
        r0.font.size = Pt(9)
        r0.font.bold = True
        r0.font.name = "Segoe UI"
        r0.font.color.rgb = tag_color
        p1 = tf.add_paragraph()
    else:
        p1 = tf.paragraphs[0]

    p1.text = title
    p1.space_after = Pt(6)
    r1 = p1.runs[0]
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.name = "Segoe UI"
    r1.font.color.rgb = SOIL

    p2 = tf.add_paragraph()
    p2.text = body
    r2 = p2.runs[0]
    r2.font.size = Pt(11)
    r2.font.name = "Segoe UI"
    r2.font.color.rgb = MUTED


def add_bullet_block(slide, left, top, width, height, lines: list[str], size: int = 14) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.name = "Segoe UI"
        run.font.color.rgb = TEXT


def add_step_pill(slide, left, top, number: str, title: str, body: str) -> None:
    w, h = Inches(1.85), Inches(1.55)
    add_card(slide, left, top, w, h, title, body, tag=f"{number}")


def ensure_notes_layout(slide, template_slide) -> None:
    """Copy notes placeholders when a new slide has an empty notes layout."""
    if slide.notes_slide.notes_text_frame is not None:
        return
    src_c_sld = template_slide.notes_slide.element.find(qn("p:cSld"))
    dst_c_sld = slide.notes_slide.element.find(qn("p:cSld"))
    if src_c_sld is None or dst_c_sld is None:
        return
    src_sp_tree = src_c_sld.find(qn("p:spTree"))
    if src_sp_tree is None:
        return
    dst_sp_tree = dst_c_sld.find(qn("p:spTree"))
    if dst_sp_tree is not None:
        dst_c_sld.remove(dst_sp_tree)
    dst_c_sld.append(deepcopy(src_sp_tree))


def set_notes(slide, text: str, *, template_slide=None) -> None:
    if template_slide is not None:
        ensure_notes_layout(slide, template_slide)
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is not None:
        text_frame.text = text
        return
    for shape in notes_slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            shape.text = text
            return


def slide_title(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_woven_accent(slide, top=Inches(0))

    add_textbox(
        slide,
        Inches(0.75),
        Inches(1.6),
        Inches(11),
        Inches(1.2),
        "Aniskwela",
        size=54,
        bold=True,
        color=SOIL,
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(2.75),
        Inches(11),
        Inches(0.8),
        "Your learning grows like a field. Slowly, visibly, and it stays yours.",
        size=22,
        color=TEXT,
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(3.65),
        Inches(11),
        Inches(0.5),
        "AI education with verifiable credentials for Filipino farmers",
        size=16,
        color=MUTED,
    )

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(4.45), Inches(2.2), Inches(0.42)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = "GRANT & PARTNER PRESENTATION"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.name = "Segoe UI"
    run.font.color.rgb = WHITE

    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.35),
        Inches(11),
        Inches(0.4),
        "TEAM AXON ENJIN",
        size=12,
        bold=True,
        color=GROWTH,
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.75),
        Inches(11),
        Inches(0.4),
        "Carlos Jerico Dela Torre  ·  Rhandie Sales Jr.  ·  Aidan Tiu",
        size=14,
        color=TEXT,
    )
    add_footer(slide, 1, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_problem(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "THE PROBLEM")
    add_headline(slide, "Rural learners lack proof, access, and a fair path to funding.")
    add_subhead(
        slide,
        "Courses assume fast Wi-Fi and English fluency. Funders cannot target or audit merit-based grants with confidence.",
    )

    cards = [
        (
            Inches(0.75),
            "English-first platforms",
            "Built for urban bandwidth. A shared 1GB phone on prepaid 3G is an afterthought.",
        ),
        (
            Inches(4.45),
            "Nothing verifiable",
            "Learners finish modules but have no portable credential employers or NGOs trust.",
        ),
        (
            Inches(8.15),
            "No grant targeting",
            "Program officers rely on forms and self-reporting instead of an auditable merit record.",
        ),
    ]
    for left, title, body in cards:
        add_card(slide, left, Inches(2.75), Inches(3.35), Inches(2.1), title, body)

    persona = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.15), Inches(11.85), Inches(1.35)
    )
    persona.fill.solid()
    persona.fill.fore_color.rgb = RGBColor(0xF3, 0xED, 0xE4)
    persona.line.color.rgb = BORDER
    tf = persona.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = (
        "Maricel (learner), Ramon (extension teacher), Divina (foundation officer): "
        "three people, one gap. Effort without proof, and funding without transparency."
    )
    p.runs[0].font.size = Pt(13)
    p.runs[0].font.name = "Segoe UI"
    p.runs[0].font.color.rgb = TEXT

    add_footer(slide, 2, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_market(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "THE OPPORTUNITY")
    add_headline(slide, "A farmer-first market that is already on mobile.")
    add_subhead(
        slide,
        "Millions of agricultural workers. Most Filipinos already online on mobile. "
        "Public skills funding exists but rarely reaches rural farmers with proof.",
    )

    cards = [
        (
            Inches(0.75),
            Inches(2.55),
            "10.7M agri workers (RSBSA)",
            "6.8M farmers · 2.7M fisherfolk · ~2M farm workers · 328K farm youth. "
            "~21% of employed Filipinos work in agriculture (PSA).",
            "TAM",
            GROWTH,
        ),
        (
            Inches(6.55),
            Inches(2.55),
            "98M online · 137M mobile",
            "83.8% internet penetration. ~89% of mobile on 3G/4G/5G. "
            "Validates EN/Fil mobile-first, sub-3G design.",
            "Reach",
            GROWTH,
        ),
        (
            Inches(0.75),
            Inches(4.55),
            "~₱16B public TVET spend",
            "TESDA free tech-voc in 2025; only ~250K of 1.37M graduates were "
            "scholarship-funded. FY2026 workforce budget ~₱18.4B.",
            "Funding pool",
            PRIMARY,
        ),
        (
            Inches(6.55),
            Inches(4.55),
            "90-day pilot beachhead",
            "50 published courses · 500 learners · 1 simulated grant program · "
            "10 extension teachers in partner network.",
            "Our wedge",
            WARNING,
        ),
    ]
    for left, top, title, body, tag, tag_color in cards:
        add_card(slide, left, top, Inches(5.55), Inches(1.75), title, body, tag=tag, tag_color=tag_color)

    add_textbox(
        slide,
        Inches(0.75),
        Inches(6.45),
        Inches(11.85),
        Inches(0.45),
        "Sources: PSA LFS 2025 · DA/ATI RSBSA (Senate hearing, Jan 2026) · "
        "DataReportal Digital 2026 · TESDA / FY2026 budget briefer",
        size=9,
        color=MUTED,
    )

    add_footer(slide, 3, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_solution(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "THE SOLUTION")
    add_headline(slide, "Aniskwela: harvest school for Filipino farmers.")
    add_subhead(
        slide,
        "A content-agnostic learning engine, positioned farmer-first, with open credentials and a funder-ready merit ledger.",
    )

    pillars = [
        (
            Inches(0.75),
            Inches(2.65),
            "AI course generation",
            "Teachers upload a PDF. Azure AI Foundry drafts modules and quizzes. Human review before publish.",
            "Live",
            GROWTH,
        ),
        (
            Inches(6.55),
            Inches(2.65),
            "Merit ledger",
            "XP, streaks, and badges accumulate. Never spent. Funders read consistency, not casino engagement.",
            "Roadmap",
            WARNING,
        ),
        (
            Inches(0.75),
            Inches(4.85),
            "Verifiable credentials",
            "W3C VC and Open Badges 3.0. Only a hash anchored on Stellar. No personal data on-chain.",
            "Roadmap",
            WARNING,
        ),
        (
            Inches(6.55),
            Inches(4.85),
            "Filipino-first, low-resource",
            "EN and Filipino toggle. Sub-3G performance budget. System fonts. Light theme only.",
            "Live",
            GROWTH,
        ),
    ]
    for left, top, title, body, tag, tag_color in pillars:
        add_card(slide, left, top, Inches(5.55), Inches(1.85), title, body, tag=tag, tag_color=tag_color)

    add_footer(slide, 4, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_learner_ui(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "LEARNER EXPERIENCE")
    add_headline(slide, "Built for the phone Maricel actually uses.")
    add_subhead(slide, "Mobile-first catalog and lessons. English and Filipino. Cached reads with no AI on the learner path.")

    # Wireframe placeholder left
    frame_l = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(2.55), Inches(3.6), Inches(4.2)
    )
    frame_l.fill.solid()
    frame_l.fill.fore_color.rgb = SURFACE
    frame_l.line.color.rgb = PRIMARY
    frame_l.line.width = Pt(2)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(2.85),
        Inches(3.1),
        Inches(0.35),
        "Course catalog",
        size=13,
        bold=True,
        color=SOIL,
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(3.35),
        Inches(3.1),
        Inches(0.8),
        "Organic Rice Basics\nFinancial Literacy 101\nSoil Health Module",
        size=11,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(1.0),
        Inches(5.9),
        Inches(3.1),
        Inches(0.5),
        "Live · drop screenshot here",
        size=10,
        color=GROWTH,
    )

    # Wireframe placeholder right
    frame_r = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.65), Inches(2.55), Inches(3.6), Inches(4.2)
    )
    frame_r.fill.solid()
    frame_r.fill.fore_color.rgb = SURFACE
    frame_r.line.color.rgb = PRIMARY
    frame_r.line.width = Pt(2)
    add_textbox(
        slide,
        Inches(4.9),
        Inches(2.85),
        Inches(3.1),
        Inches(0.35),
        "Lesson reader",
        size=13,
        bold=True,
        color=SOIL,
    )
    add_textbox(
        slide,
        Inches(4.9),
        Inches(3.35),
        Inches(3.1),
        Inches(1.2),
        "Short chunks for 3G.\nReadable body text.\nFilipino / English toggle in header.",
        size=11,
        color=MUTED,
    )
    add_textbox(
        slide,
        Inches(4.9),
        Inches(5.9),
        Inches(3.1),
        Inches(0.5),
        "Live · drop screenshot here",
        size=10,
        color=GROWTH,
    )

    bullets = [
        "≈174 KB gz shared JS (budget: 220 KB)",
        "Seed → Sprout → Scholar → Expert → Mentor levels",
        "Roadmap: quizzes, XP bar, credential wallet",
    ]
    add_bullet_block(slide, Inches(8.55), Inches(2.75), Inches(4.05), Inches(3.5), bullets, size=13)

    add_footer(slide, 5, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_teacher_funder(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "TEACHER & FUNDER SURFACES")
    add_headline(slide, "Teachers publish in an afternoon. Funders get criteria they can audit.")
    add_subhead(slide, "Upload and generate is live today. The funder eligibility console is specified and schema-ready.")

    add_card(
        slide,
        Inches(0.75),
        Inches(2.65),
        Inches(5.8),
        Inches(3.5),
        "Teacher dashboard",
        "1. Upload PDF or text\n2. Server extracts and calls AI once\n3. Review draft in dashboard\n4. Publish to public catalog\n\nNo auto-publish. Ownership enforced with RLS.",
        tag="Live",
        tag_color=GROWTH,
    )
    add_card(
        slide,
        Inches(6.85),
        Inches(2.65),
        Inches(5.75),
        Inches(3.5),
        "Funder program builder",
        "Define criteria against the merit ledger.\nPreview eligible learners.\nRun simulated disbursement (labelled).\nExport audit report for donors.\n\nReal payout via licensed VASP partner.",
        tag="Roadmap",
        tag_color=WARNING,
    )

    add_footer(slide, 6, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_how_it_works(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "HOW IT WORKS")
    add_headline(slide, "From document to anchored credential.")
    add_subhead(slide, "AI runs once per course. Learners never wait on a model when they open a lesson.")

    steps = [
        ("1", "Upload", "Teacher sends PDF or doc."),
        ("2", "Generate", "Azure AI Foundry drafts structure."),
        ("3", "Review", "Teacher edits and publishes."),
        ("4", "Learn", "Learner completes lessons and quizzes."),
        ("5", "Credential", "Platform issues signed W3C VC."),
        ("6", "Anchor", "Hash recorded on Stellar Testnet."),
    ]
    x_start = Inches(0.55)
    for i, (num, title, body) in enumerate(steps):
        add_step_pill(slide, x_start + Inches(i * 2.05), Inches(2.85), num, title, body)

    add_card(
        slide,
        Inches(0.75),
        Inches(5.0),
        Inches(5.5),
        Inches(1.35),
        "Human in the loop",
        "Mandatory teacher review before any course goes live.",
    )
    add_card(
        slide,
        Inches(6.55),
        Inches(5.0),
        Inches(5.5),
        Inches(1.35),
        "Demo resilience",
        "If Testnet is down, a labelled mock anchor keeps the demo moving.",
    )

    add_footer(slide, 7, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_credentials(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "TRUST LAYER")
    add_headline(slide, "Merit ledger plus verifiable credentials.")
    add_subhead(slide, "Open standards. Hash-only on-chain. A public verifier anyone can use without an account.")

    add_card(
        slide,
        Inches(0.75),
        Inches(2.65),
        Inches(3.7),
        Inches(2.4),
        "Merit ledger",
        "Append-only XP events. Badges for consistency. Readable by funders when setting eligibility rules.",
        tag="Roadmap",
        tag_color=WARNING,
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(2.65),
        Inches(3.7),
        Inches(2.4),
        "W3C VC + OB 3.0",
        "Portable credential any OB 3.0 verifier can check. Learner name on the card, not on the chain.",
        tag="Roadmap",
        tag_color=WARNING,
    )
    add_card(
        slide,
        Inches(8.75),
        Inches(2.65),
        Inches(3.85),
        Inches(2.4),
        "Public verifier",
        "URL or QR → valid or invalid, with course, date, and score. No login required.",
        tag="Roadmap",
        tag_color=WARNING,
    )

    add_card(
        slide,
        Inches(0.75),
        Inches(5.35),
        Inches(11.85),
        Inches(1.15),
        "Regulatory posture",
        "Aniskwela is not a money transmitter. Credentials carry proof of learning. Grants use simulated disbursement until a VASP partner executes payout.",
    )

    add_footer(slide, 8, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_grant_model(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "FOR FUNDERS & PARTNERS")
    add_headline(slide, "Target grants to consistent learners. Audit every decision.")
    add_subhead(slide, "Aniskwela decides eligibility. A licensed VASP moves money. You get exports your donors can read.")

    add_card(
        slide,
        Inches(0.75),
        Inches(2.65),
        Inches(3.7),
        Inches(2.5),
        "Set criteria",
        'Example: "Agriculture track, ≥30,000 XP, Consistent Learner badge." Criteria stored as structured JSON.',
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(2.65),
        Inches(3.7),
        Inches(2.5),
        "Preview eligible list",
        "Server-side query against merit ledger and credentials. No self-reported forms.",
    )
    add_card(
        slide,
        Inches(8.75),
        Inches(2.65),
        Inches(3.85),
        Inches(2.5),
        "Simulate & export",
        "MVP disbursement is clearly labelled simulation. Audit report export for program officers.",
        tag="Pilot ask",
        tag_color=PRIMARY,
    )

    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.45),
        Inches(11.85),
        Inches(0.9),
        "We are looking for one foundation or agency partner to co-design the first pilot program criteria with real farmers and extension teachers in your network.",
        size=15,
        color=TEXT,
    )

    add_footer(slide, 9, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_tech(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "TECHNOLOGY")
    add_headline(slide, "Honest stack for rural scale.")
    add_subhead(slide, "Next.js 16 on Vercel. Supabase with RLS on every table. Azure AI Foundry. Stellar for hash anchoring only.")

    stack = [
        ("Framework", "Next.js 16.2 App Router", "Cache Components on public reads"),
        ("Data", "Supabase Postgres", "9 tables, RLS, auth SSR"),
        ("AI", "Azure AI Foundry", "gpt-5.4 course gen · gpt-5.4-mini repair"),
        ("Chain", "Stellar Testnet", "Hash-only memo · mock fallback flag"),
    ]
    y = Inches(2.65)
    for label, title, detail in stack:
        add_card(slide, Inches(0.75), y, Inches(5.5), Inches(0.95), f"{label}: {title}", detail)
        y += Inches(1.05)

    perf_lines = [
        "Initial JS ≤ 220 KB gzipped (measured ≈174 KB today)",
        "Images ≤ 80 KB WebP on first load",
        "Core content target < 5 s on 3G",
        "System fonts only · light theme · EN + Filipino",
    ]
    add_bullet_block(slide, Inches(6.85), Inches(2.75), Inches(5.75), Inches(3.5), perf_lines, size=13)

    add_footer(slide, 10, TOTAL_SLIDES)
    set_notes(slide, notes)


def slide_cta(prs, notes: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_kicker(slide, "NEXT STEP")
    add_headline(slide, "Help us pilot with real farmers.")
    add_subhead(slide, "Co-design eligibility rules. Connect ten extension teachers. Give feedback when the verifier ships.")

    metrics = [
        ("50", "published courses\n(90 days)"),
        ("500", "registered learners\n(90 days)"),
        ("1", "simulated grant program\nwith partner audit export"),
    ]
    x = Inches(0.75)
    for num, label in metrics:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(2.75), Inches(3.5), Inches(1.65)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = SURFACE
        box.line.color.rgb = GROWTH
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = num
        r1.font.size = Pt(36)
        r1.font.bold = True
        r1.font.color.rgb = GROWTH
        r1.font.name = "Segoe UI"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED
        r2.font.name = "Segoe UI"
        x += Inches(3.85)

    qr = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(4.85), Inches(1.35), Inches(1.35)
    )
    qr.fill.solid()
    qr.fill.fore_color.rgb = SURFACE
    qr.line.color.rgb = BORDER
    add_textbox(
        slide,
        Inches(0.85),
        Inches(5.25),
        Inches(1.15),
        Inches(0.6),
        "QR",
        size=14,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(2.35),
        Inches(4.95),
        Inches(6),
        Inches(1.2),
        "Try the live app\n(link your Vercel URL before presenting)",
        size=16,
        color=TEXT,
    )

    add_textbox(
        slide,
        Inches(0.75),
        Inches(6.35),
        Inches(11),
        Inches(0.5),
        "Team Axon Enjin  ·  Dela Torre  ·  Sales Jr.  ·  Tiu",
        size=14,
        bold=True,
        color=SOIL,
    )

    add_footer(slide, 11, TOTAL_SLIDES)
    set_notes(slide, notes)


NOTES = [
    """Good morning. We are Team Axon Enjin, and this is Aniskwela.

The name comes from two Filipino words: ani, harvest, and eskwela, school. We built it for Filipino farmers and rural learners who study on shared phones and prepaid data.

Aniskwela turns raw teaching materials into structured courses, helps learners build a visible record of effort, and issues credentials that funders can verify without taking anyone's word for it.

Today we want to show you why that matters for the programs you run, what we have built so far, and what a partnership could look like.""",
    """Let me start with Maricel. She is twenty-two, lives in a rural part of Region IV, and shares a one-gigabyte Android phone with her family. She wants to learn agriculture and basic financial literacy, but the platforms she finds were not built for her life.

The courses are mostly in English. They assume fast Wi-Fi. When she finishes something, she has nothing credible to show an employer or a grant officer.

Teachers face the opposite problem. Ramon is an agricultural extension worker. He knows the subject, but building an online course takes tools and time he does not have.

Funders like you see the third gap. Divina runs a foundation program. She wants to reward consistent learners, but she cannot target grants transparently and she cannot audit outcomes without manual paperwork.

Existing tools either gamify rewards without real learning, or they offer credentials that are easy to fake. Neither side gets trust.""",
    """Before we show the product, let me quantify the opportunity.

The Registry System for Basic Sectors in Agriculture lists ten point seven million agricultural workers—nearly seven million farmers, two point seven million fisherfolk, plus farm workers and farm youth. Agriculture still employs about one in five Filipino workers according to PSA.

At the same time, ninety-eight million Filipinos are online—eighty-four percent penetration—and most mobile connections run on 3G, 4G, or 5G. That is why we built for prepaid data and shared phones, not campus Wi-Fi.

Public money for skills training already exists. TESDA spent roughly sixteen billion pesos on free tech-voc in twenty twenty-five, but fewer than one quarter of graduates went through scholarship vouchers. The funding pool is there; what is missing is transparent, merit-based targeting for rural farmers.

Our ninety-day wedge is modest: fifty courses, five hundred learners, one simulated grant program, and ten extension teachers through a partner network. That is what we are asking you to help us pilot.""",
    """Aniskwela addresses all three sides with one platform.

First, teachers upload a PDF or document and our AI pipeline drafts a structured course. Nothing goes live until a teacher reviews it. We do not auto-publish.

Second, learners study on a low-resource interface tuned for 3G and Filipino. Progress shows up as XP and levels, from Seed to Mentor. That record is cumulative. We do not burn or spend XP like a casino game.

Third, when a learner passes a course, we issue a W3C Verifiable Credential in the Open Badges 3.0 profile. Only a hash of that credential goes on Stellar. No personal data on the chain.

Fourth, funders define eligibility against the merit ledger. Aniskwela decides who qualifies. A licensed VASP moves money. We stay on the right side of regulation.

The engine can serve other subjects later, but we are farmer-first because that is where the need and our team's story align.""",
    """This is what a learner sees today. We already ship a course catalog and lesson reader that load quickly on a typical Philippine phone screen.

The shared initial JavaScript bundle measures about one hundred seventy-four kilobytes gzipped, under our two hundred twenty kilobyte budget. Content is cached on the server so learners are not waiting on AI when they open a lesson.

Language toggles between English and Filipino without a full page reload. That matters when someone reads faster in Filipino but still wants English for certain terms.

What is live now: browsing published courses, reading lessons, switching language.

What is on the roadmap for this surface: interactive quizzes, XP and streak display, enrollments, and the credential wallet.

We label that honestly so you know what you can click today versus what we are building next.""",
    """On the teacher side, the flow that works today is upload, generate, review, and publish.

A teacher signs in, uploads a PDF, and the server extracts text and calls Azure AI Foundry once per course. The model returns modules, lessons, and quiz questions. If the JSON fails validation, a smaller model tries one repair pass. If that still fails, we return an error. We never save garbage.

The teacher sees the draft in their dashboard and publishes when ready. Published courses appear in the public catalog.

The funder dashboard is not built yet, but the database already has grant program tables with a simulated flag. The design lets a program officer set criteria like agriculture learners above a certain XP threshold, preview the eligible list, and run a labelled simulation with an exportable audit trail. Real disbursement always routes through a licensed partner.""",
    """Here is the end-to-end story we are completing.

Step one: a teacher uploads source material.

Step two: Azure AI Foundry generates a structured draft. AI runs once per course, not on every page view.

Step three: the teacher edits if needed and publishes.

Step four: a learner studies lessons and takes assessments. We record completion in the merit ledger.

Step five: on pass, the platform issues a verifiable credential and signs it with our issuer key.

Step six: we anchor only the credential hash in a Stellar transaction. Anyone can verify the link between the credential and the chain entry without seeing private learner data on ledger.

If Testnet is down during a demo, we fall back to a clearly labelled mock anchor so learning never blocks on infrastructure.""",
    """This slide is about trust without hype.

The merit ledger is an append-only record of positive events: lesson completed, quiz passed, badge earned. XP never goes negative and never gets spent. Funders read that ledger when they ask who stayed consistent.

Credentials follow open standards, not a proprietary badge NFT. Employers and other institutions can verify them with tools they already trust.

On-chain, we store a hash. That proves the credential existed at issuance time and was not altered afterward. We do not put names, emails, or grades in the memo.

A public verifier page, no login required, will let anyone scan a QR or open a link and see valid or invalid with course, date, and score. That piece is specified and schema-ready; we are building it in the next sprint.""",
    """This is the slide most relevant to your work.

Aniskwela is not a money transmitter. We do not move pesos inside the app. We tell you who meets the criteria you set. Your licensed VASP or e-money partner executes payout.

In the MVP, disbursement is a simulation. Every screen and export says so clearly. That lets you pilot policy, report to donors, and train staff before real funds flow.

You might fund learners who completed an agriculture track, held a thirty-day streak, and earned a specific badge. The platform returns a list backed by ledger data, not self-reported forms.

For compliance, we lean on row-level security in Postgres, server-side eligibility checks, and exportable audit logs. CLR work covers privacy and trademark next steps.

We are looking for one foundation or agency partner willing to co-design the first pilot program criteria.""",
    """Under the hood we use a stack chosen for speed, cost, and honesty about constraints.

The app is Next.js sixteen on Vercel. Auth, database, and storage run on Supabase with row-level security on every table. Teachers only see their drafts; learners only see their own progress rows when those flows ship.

AI goes through Azure AI Foundry with GPT five point four for course generation and the mini model for cheap repair tasks. Managed Identity in production, API key only for local development.

The initial migration defines nine tables: profiles, courses, lessons, quiz questions, enrollments, merit ledger, badges, credentials, and grant programs. Several are wired in schema today and waiting for application code.

Performance is a requirement, not a nice-to-have. Light theme, system fonts, no web font download on first paint, images capped at eighty kilobytes WebP where we use them.""",
    """We are not asking you to fund a slide deck. We are asking for a conversation about a pilot.

Help us define real eligibility rules for farmers in your network. Point us to ten teachers who can upload extension materials. Give us feedback when the verifier and funder console land in the next build weeks.

Success for us in the next ninety days looks like fifty published courses, five hundred registered learners, and at least one simulated grant program run with a partner who exports an audit report they would actually show a donor.

Scan the QR or open the link to try the live teacher and catalog flows. Talk to us after about co-designing the first grant criteria.

Thank you. We are Carlos, Rhandie, and Aidan from Team Axon Enjin. Aniskwela: learn, earn credentials you can prove, and grow.""",
]


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs, NOTES[0])
    slide_problem(prs, NOTES[1])
    slide_market(prs, NOTES[2])
    slide_solution(prs, NOTES[3])
    slide_learner_ui(prs, NOTES[4])
    slide_teacher_funder(prs, NOTES[5])
    slide_how_it_works(prs, NOTES[6])
    slide_credentials(prs, NOTES[7])
    slide_grant_model(prs, NOTES[8])
    slide_tech(prs, NOTES[9])
    slide_cta(prs, NOTES[10])

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path} ({path.stat().st_size:,} bytes)")
