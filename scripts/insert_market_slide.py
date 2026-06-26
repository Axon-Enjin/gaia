#!/usr/bin/env python3
"""
Insert the Market / Opportunity slide into an existing Aniskwela-Grant-Pitch.pptx
without regenerating the full deck (preserves manual edits).

Usage (from repo root):
    python scripts/insert_market_slide.py
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation

REPO_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO_ROOT / "scripts"))

from build_pitch_deck import (  # noqa: E402
    DECK_LABEL,
    NOTES,
    OUTPUT,
    TOTAL_SLIDES,
    set_notes,
    slide_market,
)

INSERT_AT = 2  # zero-based index → slide 3


def slide_has_kicker(slide, kicker: str) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if kicker in para.text:
                return True
    return False


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    slide_id = sld_id_lst[old_index]
    sld_id_lst.remove(slide_id)
    sld_id_lst.insert(new_index, slide_id)


def market_slide_index(prs: Presentation) -> int | None:
    for i, slide in enumerate(prs.slides):
        if slide_has_kicker(slide, "THE OPPORTUNITY"):
            return i
    return None


def renumber_footers(prs: Presentation, total: int = TOTAL_SLIDES) -> None:
    footer_pattern = re.compile(r"Grant partner deck\s*.\s*\d+/\d+")
    for slide_index, slide in enumerate(prs.slides):
        page = slide_index + 1
        replacement = f"{DECK_LABEL}  ·  {page}/{total}"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if "Grant partner deck" not in para.text:
                    continue
                if footer_pattern.search(para.text):
                    para.text = replacement


def insert_market_slide(path: Path = OUTPUT) -> Path:
    prs = Presentation(str(path))
    existing = market_slide_index(prs)

    if existing is None:
        slide_market(prs, NOTES[2])
        existing = len(prs.slides) - 1

    if existing != INSERT_AT:
        move_slide(prs, existing, INSERT_AT)

    market_slide = prs.slides[INSERT_AT]
    template = prs.slides[1] if len(prs.slides) > 1 else prs.slides[0]
    set_notes(market_slide, NOTES[2], template_slide=template)
    renumber_footers(prs)
    prs.save(str(path))
    return path


if __name__ == "__main__":
    out = insert_market_slide()
    print(f"Updated {out} ({out.stat().st_size:,} bytes, {TOTAL_SLIDES} slides)")
